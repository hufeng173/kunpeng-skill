#!/usr/bin/env python3
"""Extract and structure local documents for writing and method distillation."""

from __future__ import annotations

import argparse
import hashlib
import json
import math
import re
from collections import Counter
from pathlib import Path
from typing import Any, Iterable

from kunpeng_common import (
    DOCUMENT_EXTENSIONS,
    aggregate_status,
    atomic_write_json,
    atomic_write_text,
    bounded_error,
    configure_utf8,
    find_sources,
    prepare_output,
    quantile,
    relative_artifact,
    reused_analysis_status,
    sampled_fingerprint,
    source_id,
    status_counts,
    utc_now,
)


SENTENCE_RE = re.compile(r"[^。！？!?\.]+[。！？!?\.]?|[^。！？!?\.]+$", re.MULTILINE)
ASCII_WORD_RE = re.compile(r"[A-Za-z0-9][A-Za-z0-9'_-]*")
CJK_RE = re.compile(r"[\u3400-\u4dbf\u4e00-\u9fff]")
HEADING_RE = re.compile(r"(?m)^\s{0,3}(#{1,6})\s+(.+?)\s*$")
LIST_RE = re.compile(r"(?m)^\s*(?:[-*+] |\d+[.)、]\s*)")
TIMESTAMP_LINE_RE = re.compile(
    r"^\s*(?:\d+:)?\d{1,2}:\d{2}[,.]\d{3}\s*-->\s*(?:\d+:)?\d{1,2}:\d{2}[,.]\d{3}.*$"
)
TAG_RE = re.compile(r"<[^>]+>|\{\\[^}]+\}")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Extract local documents and compute reproducible writing-style evidence."
    )
    parser.add_argument("source", type=Path, help="Document file or directory.")
    parser.add_argument("--output", type=Path, required=True)
    parser.add_argument("--no-recursive", action="store_true")
    parser.add_argument("--resume", action="store_true")
    parser.add_argument("--ocr-scanned", choices=("auto", "on", "off"), default="auto")
    parser.add_argument("--ocr-lang", default="ch")
    parser.add_argument("--ocr-device", default="auto")
    parser.add_argument("--chunk-chars", type=int, default=6000)
    parser.add_argument("--max-documents", type=int, default=500)
    parser.add_argument("--max-pdf-pages", type=int, default=2000)
    return parser.parse_args()


def read_text_file(path: Path) -> str:
    data = path.read_bytes()
    for encoding in ("utf-8-sig", "utf-16", "gb18030", "cp1252"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def extract_plain(path: Path) -> tuple[str, list[dict[str, Any]], dict[str, Any]]:
    text = read_text_file(path).replace("\r\n", "\n")
    outline = [
        {"level": len(match.group(1)), "title": match.group(2).strip()}
        for match in HEADING_RE.finditer(text)
    ]
    return text, outline[:300], {}


def extract_html(path: Path) -> tuple[str, list[dict[str, Any]], dict[str, Any]]:
    raw = read_text_file(path)
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(raw, "html.parser")
    outline = [
        {"level": int(node.name[1]), "title": node.get_text(" ", strip=True)}
        for node in soup.find_all(re.compile(r"^h[1-6]$"))
        if node.get_text(" ", strip=True)
    ]
    extractor = "beautifulsoup4"
    try:
        import trafilatura

        text = trafilatura.extract(
            raw,
            include_comments=False,
            include_tables=True,
            output_format="txt",
        ) or ""
        if text.strip():
            extractor = "trafilatura"
        else:
            text = soup.get_text("\n", strip=True)
    except ImportError:
        text = soup.get_text("\n", strip=True)
    return text, outline[:300], {"extractor": extractor}


def extract_docx(path: Path) -> tuple[str, list[dict[str, Any]], dict[str, Any]]:
    from docx import Document

    document = Document(path)
    blocks: list[str] = []
    outline: list[dict[str, Any]] = []
    for paragraph in document.paragraphs:
        value = paragraph.text.strip()
        if not value:
            continue
        blocks.append(value)
        style = paragraph.style.name if paragraph.style else ""
        match = re.match(r"Heading\s+(\d+)", style, re.IGNORECASE)
        if match:
            outline.append({"level": int(match.group(1)), "title": value})
    for table_index, table in enumerate(document.tables, start=1):
        blocks.append(f"[Table {table_index}]")
        for row in table.rows:
            blocks.append("\t".join(cell.text.strip() for cell in row.cells))
    return "\n\n".join(blocks), outline[:300], {"table_count": len(document.tables)}


def extract_pptx(path: Path) -> tuple[str, list[dict[str, Any]], dict[str, Any]]:
    from pptx import Presentation

    presentation = Presentation(path)
    blocks: list[str] = []
    outline: list[dict[str, Any]] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            value = getattr(shape, "text", "").strip()
            if value:
                texts.append(value)
        title = texts[0].splitlines()[0][:200] if texts else f"Slide {slide_index}"
        outline.append({"level": 1, "title": title, "slide": slide_index})
        blocks.append(f"[Slide {slide_index}]\n" + "\n".join(texts))
    return "\n\n".join(blocks), outline, {"slide_count": len(presentation.slides)}


def extract_subtitles(path: Path) -> tuple[str, list[dict[str, Any]], dict[str, Any]]:
    try:
        import pysubs2

        subtitles = pysubs2.load(str(path), encoding="utf-8")
        lines = [event.plaintext.strip() for event in subtitles if event.plaintext.strip()]
        return "\n".join(lines), [], {"subtitle_events": len(lines), "extractor": "pysubs2"}
    except (ImportError, UnicodeError, OSError):
        lines = []
        for raw_line in read_text_file(path).splitlines():
            line = raw_line.strip()
            if not line or line.isdigit() or TIMESTAMP_LINE_RE.match(line):
                continue
            cleaned = TAG_RE.sub("", line).strip()
            if cleaned:
                lines.append(cleaned)
        return "\n".join(lines), [], {"subtitle_events": len(lines), "extractor": "fallback"}


def extract_pdf(
    path: Path,
    item_dir: Path,
    ocr_engine: Any | None,
    max_pages: int,
) -> tuple[str, list[dict[str, Any]], dict[str, Any]]:
    from pypdf import PdfReader

    reader = PdfReader(path)
    if len(reader.pages) > max_pages:
        raise RuntimeError(f"PDF has {len(reader.pages)} pages; raise --max-pdf-pages to process it.")
    page_texts: list[str] = []
    low_text_pages: list[int] = []
    for index, page in enumerate(reader.pages):
        value = (page.extract_text() or "").strip()
        page_texts.append(value)
        if len(value) < 30:
            low_text_pages.append(index)

    ocr_pages: list[int] = []
    rendered_pages: list[dict[str, Any]] = []
    ocr_errors: list[str] = []
    if low_text_pages:
        try:
            import pypdfium2 as pdfium

            pdf = pdfium.PdfDocument(path)
            render_dir = item_dir / "ocr-pages"
            render_dir.mkdir(parents=True, exist_ok=True)
            for index in low_text_pages:
                image_path = render_dir / f"page-{index + 1:04d}.jpg"
                try:
                    bitmap = pdf[index].render(scale=2.0)
                    bitmap.to_pil().convert("RGB").save(image_path, quality=90)
                    rendered_pages.append(
                        {
                            "page": index + 1,
                            "file": relative_artifact(image_path, item_dir),
                        }
                    )
                    if ocr_engine:
                        lines = ocr_engine.recognize(image_path)
                        extracted = "\n".join(
                            line["text"] for line in lines if line.get("text")
                        )
                        if extracted.strip():
                            page_texts[index] = extracted
                            ocr_pages.append(index + 1)
                except Exception as exc:
                    ocr_errors.append(f"page {index + 1}: {bounded_error(exc, path, item_dir)}")
        except ImportError:
            ocr_errors.append("pypdfium2 is not installed; low-text pages were not rendered")

    blocks = [f"[Page {index + 1}]\n{text}" for index, text in enumerate(page_texts) if text]
    unresolved_pages = [page for page in (index + 1 for index in low_text_pages) if page not in ocr_pages]
    metadata = {
        "page_count": len(reader.pages),
        "low_text_pages": [index + 1 for index in low_text_pages],
        "ocr_pages": ocr_pages,
        "rendered_pages": rendered_pages,
        "unresolved_low_text_pages": unresolved_pages,
        "ocr_errors": ocr_errors[:20],
    }
    return "\n\n".join(blocks), [], metadata


def extract_document(
    path: Path,
    item_dir: Path,
    ocr_engine: Any | None,
    max_pdf_pages: int,
) -> tuple[str, list[dict[str, Any]], dict[str, Any]]:
    suffix = path.suffix.casefold()
    if suffix in {".txt", ".text", ".md", ".rst", ".csv", ".json", ".yaml", ".yml"}:
        return extract_plain(path)
    if suffix in {".html", ".htm"}:
        return extract_html(path)
    if suffix == ".docx":
        return extract_docx(path)
    if suffix == ".pptx":
        return extract_pptx(path)
    if suffix == ".pdf":
        return extract_pdf(path, item_dir, ocr_engine, max_pdf_pages)
    if suffix in {".srt", ".vtt", ".ass"}:
        return extract_subtitles(path)
    raise RuntimeError(f"Unsupported document type: {suffix}")


def tokenize(text: str) -> list[str]:
    ascii_tokens = [token.casefold() for token in ASCII_WORD_RE.findall(text)]
    cjk_text = "".join(CJK_RE.findall(text))
    if not cjk_text:
        return ascii_tokens
    try:
        import jieba

        cjk_tokens = [token for token in jieba.lcut(cjk_text) if len(token.strip()) > 1]
    except ImportError:
        cjk_tokens = [cjk_text[index : index + 2] for index in range(max(0, len(cjk_text) - 1))]
    return ascii_tokens + cjk_tokens


def style_metrics(text: str) -> dict[str, Any]:
    lines = [line.rstrip() for line in text.splitlines()]
    nonempty = [line.strip() for line in lines if line.strip()]
    paragraphs = [part.strip() for part in re.split(r"\n\s*\n", text) if part.strip()]
    sentences = [match.group(0).strip() for match in SENTENCE_RE.finditer(text) if match.group(0).strip()]
    sentence_lengths = [len(CJK_RE.findall(value)) + len(ASCII_WORD_RE.findall(value)) for value in sentences]
    paragraph_lengths = [len(value) for value in paragraphs]
    tokens = tokenize(text)
    token_counts = Counter(tokens)
    punctuation = Counter(character for character in text if character in "，。！？；：、,.!?;:—-…（）()“”\"'")
    cjk_count = len(CJK_RE.findall(text))
    ascii_word_count = len(ASCII_WORD_RE.findall(text))
    denominator = max(1, len(nonempty))
    return {
        "characters": len(text),
        "cjk_characters": cjk_count,
        "ascii_words": ascii_word_count,
        "paragraphs": len(paragraphs),
        "sentences": len(sentences),
        "sentence_length": {
            "mean": round(sum(sentence_lengths) / max(1, len(sentence_lengths)), 2),
            "p25": round(quantile(sentence_lengths, 0.25), 2),
            "median": round(quantile(sentence_lengths, 0.50), 2),
            "p75": round(quantile(sentence_lengths, 0.75), 2),
            "max": max(sentence_lengths, default=0),
        },
        "paragraph_length": {
            "mean_chars": round(sum(paragraph_lengths) / max(1, len(paragraph_lengths)), 2),
            "p25_chars": round(quantile(paragraph_lengths, 0.25), 2),
            "median_chars": round(quantile(paragraph_lengths, 0.50), 2),
            "p75_chars": round(quantile(paragraph_lengths, 0.75), 2),
        },
        "lexical_diversity": round(len(token_counts) / max(1, len(tokens)), 4),
        "question_sentence_share": round(sum(value.endswith(("?", "？")) for value in sentences) / max(1, len(sentences)), 4),
        "exclamation_sentence_share": round(sum(value.endswith(("!", "！")) for value in sentences) / max(1, len(sentences)), 4),
        "heading_line_share": round(sum(bool(HEADING_RE.match(line)) for line in nonempty) / denominator, 4),
        "list_line_share": round(sum(bool(LIST_RE.match(line)) for line in nonempty) / denominator, 4),
        "dialogue_or_quote_marks": sum(text.count(mark) for mark in ("“", "”", "\"", "‘", "’")),
        "punctuation": dict(punctuation.most_common()),
        "top_tokens": [{"token": token, "count": count} for token, count in token_counts.most_common(30)],
    }


def split_long_block(block: str, limit: int) -> Iterable[str]:
    if len(block) <= limit:
        yield block
        return
    sentences = [match.group(0).strip() for match in SENTENCE_RE.finditer(block) if match.group(0).strip()]
    if len(sentences) <= 1:
        for start in range(0, len(block), limit):
            yield block[start : start + limit]
        return
    current: list[str] = []
    size = 0
    for sentence in sentences:
        if current and size + len(sentence) + 1 > limit:
            yield "\n".join(current)
            current, size = [], 0
        current.append(sentence)
        size += len(sentence) + 1
    if current:
        yield "\n".join(current)


def natural_chunks(text: str, limit: int) -> list[str]:
    limit = max(1000, limit)
    blocks = [part.strip() for part in re.split(r"\n\s*\n", text) if part.strip()]
    chunks: list[str] = []
    current: list[str] = []
    size = 0
    for raw_block in blocks:
        for block in split_long_block(raw_block, limit):
            if current and size + len(block) + 2 > limit:
                chunks.append("\n\n".join(current))
                current, size = [], 0
            current.append(block)
            size += len(block) + 2
    if current:
        chunks.append("\n\n".join(current))
    return chunks


def sample_text(text: str, size: int = 900) -> dict[str, str]:
    if len(text) <= size * 3:
        return {"opening": text}
    middle = max(0, len(text) // 2 - size // 2)
    return {
        "opening": text[:size],
        "middle": text[middle : middle + size],
        "ending": text[-size:],
    }


def label_for(path: Path, root: Path) -> str:
    return path.relative_to(root).as_posix() if root.is_dir() else path.name


def normalized_signature(text: str) -> str:
    compact = re.sub(r"\s+", "", text).casefold()
    return hashlib.sha256(compact.encode("utf-8")).hexdigest()


def simhash64(text: str) -> int:
    compact = re.sub(r"\s+", "", text).casefold()
    shingles = [compact[index : index + 5] for index in range(max(1, len(compact) - 4))]
    if not shingles:
        return 0
    vector = [0] * 64
    step = max(1, len(shingles) // 12000)
    for shingle in shingles[::step]:
        value = int.from_bytes(hashlib.blake2b(shingle.encode("utf-8"), digest_size=8).digest(), "big")
        for bit in range(64):
            vector[bit] += 1 if value & (1 << bit) else -1
    result = 0
    for bit, score in enumerate(vector):
        if score >= 0:
            result |= 1 << bit
    return result


def build_corpus_analysis(output: Path, items: list[dict[str, Any]]) -> dict[str, Any]:
    records: list[dict[str, Any]] = []
    line_documents: Counter[str] = Counter()
    for item in items:
        analysis_ref = item.get("analysis")
        if item.get("status") == "failed" or not analysis_ref:
            continue
        try:
            analysis = json.loads((output / analysis_ref).read_text(encoding="utf-8"))
            content_ref = analysis.get("artifacts", {}).get("content")
            if not content_ref:
                continue
            text = (output / content_ref).read_text(encoding="utf-8-sig", errors="replace")
        except (OSError, UnicodeError, json.JSONDecodeError):
            continue
        lines = {
            re.sub(r"\s+", " ", line).strip()
            for line in text.splitlines()
            if 8 <= len(re.sub(r"\s+", "", line)) <= 240
        }
        line_documents.update(lines)
        records.append(
            {
                "id": item["id"],
                "source": item.get("source"),
                "characters": len(text),
                "signature": normalized_signature(text),
                "simhash": simhash64(text),
                "metrics": analysis.get("style_metrics") or {},
            }
        )

    exact_groups: dict[str, list[str]] = {}
    for record in records:
        exact_groups.setdefault(record["signature"], []).append(record["id"])
    exact_duplicates = [group for group in exact_groups.values() if len(group) > 1]
    near_duplicates: list[dict[str, Any]] = []
    for left_index, left in enumerate(records):
        for right in records[left_index + 1 :]:
            if left["signature"] == right["signature"]:
                continue
            distance = (left["simhash"] ^ right["simhash"]).bit_count()
            if distance <= 5:
                near_duplicates.append({"left": left["id"], "right": right["id"], "simhash_distance": distance})

    def distribution(path: tuple[str, ...]) -> dict[str, float]:
        values: list[float] = []
        for record in records:
            value: Any = record["metrics"]
            for key in path:
                value = value.get(key) if isinstance(value, dict) else None
            if isinstance(value, (int, float)):
                values.append(float(value))
        return {
            "mean": round(sum(values) / max(1, len(values)), 4),
            "p25": round(quantile(values, 0.25), 4),
            "median": round(quantile(values, 0.5), 4),
            "p75": round(quantile(values, 0.75), 4),
        }

    repeated_threshold = max(2, math.ceil(len(records) * 0.6)) if records else 2
    repeated_lines = [
        {"text": line[:240], "document_count": count, "document_share": round(count / max(1, len(records)), 4)}
        for line, count in line_documents.most_common(100)
        if count >= repeated_threshold
    ]
    return {
        "schema_version": 1,
        "source_count": len(records),
        "total_characters": sum(record["characters"] for record in records),
        "exact_duplicate_groups": exact_duplicates,
        "near_duplicate_candidates": near_duplicates[:1000],
        "repeated_line_candidates": repeated_lines,
        "style_distributions": {
            "sentence_length_mean": distribution(("sentence_length", "mean")),
            "paragraph_length_mean": distribution(("paragraph_length", "mean_chars")),
            "lexical_diversity": distribution(("lexical_diversity",)),
            "question_sentence_share": distribution(("question_sentence_share",)),
            "exclamation_sentence_share": distribution(("exclamation_sentence_share",)),
            "heading_line_share": distribution(("heading_line_share",)),
            "list_line_share": distribution(("list_line_share",)),
        },
        "distillation_status": "evidence_ready",
        "semantic_review_required": [
            "remove boilerplate, duplicates, reprints, and mixed-author outliers before aggregation",
            "create one evidence-linked semantic card per retained source",
            "separate topic-dependent vocabulary from stable structure, voice, reasoning, teaching, and rhetorical mechanisms",
        ],
    }


def main() -> int:
    configure_utf8()
    args = parse_args()
    sources = find_sources(args.source, DOCUMENT_EXTENSIONS, not args.no_recursive)
    if not sources:
        raise SystemExit("No supported documents found.")
    if len(sources) > max(1, args.max_documents):
        raise SystemExit(
            f"Found {len(sources)} documents; raise --max-documents to process them all."
        )
    output = prepare_output(args.output, args.resume)
    source_root = args.source.resolve()

    ocr_engine = None
    ocr_error = None
    if args.ocr_scanned != "off" and any(path.suffix.casefold() == ".pdf" for path in sources):
        try:
            from local_ocr import LocalOCR

            ocr_engine = LocalOCR(args.ocr_lang, args.ocr_device)
        except Exception as exc:
            ocr_error = bounded_error(exc)

    items: list[dict[str, Any]] = []
    for path in sources:
        item_id = source_id(path)
        item_dir = output / "documents" / item_id
        analysis_path = item_dir / "analysis.json"
        label = label_for(path, source_root)
        if args.resume and analysis_path.exists():
            reused_status = reused_analysis_status(analysis_path)
            items.append(
                {
                    "id": item_id,
                    "source": label,
                    "status": reused_status,
                    "extraction_status": reused_status,
                    "distillation_status": "evidence_ready",
                    "reused": True,
                    "analysis": relative_artifact(analysis_path, output),
                }
            )
            continue
        try:
            text, outline, extraction = extract_document(
                path, item_dir, ocr_engine, max(1, args.max_pdf_pages)
            )
            normalized = text.replace("\r\n", "\n").strip()
            is_pdf = path.suffix.casefold() == ".pdf"
            rendered_pages = extraction.get("rendered_pages", []) if is_pdf else []
            unresolved_pages = extraction.get("unresolved_low_text_pages", []) if is_pdf else []
            if not normalized and not rendered_pages:
                raise RuntimeError("No text could be extracted and no pages could be rendered")

            stages: dict[str, dict[str, Any]] = {
                "text_extraction": {
                    "status": "complete" if normalized else "partial",
                    **(
                        {
                            "fallback": "host_visual_review",
                            "host_review_required": True,
                        }
                        if not normalized
                        else {}
                    ),
                }
            }
            if is_pdf and extraction.get("low_text_pages"):
                if not unresolved_pages:
                    scanned_status = "complete"
                else:
                    scanned_status = "partial"
                stages["scanned_pdf_text"] = {
                    "status": scanned_status,
                    "unresolved_pages": unresolved_pages,
                    "fallback": "host_visual_review" if rendered_pages else None,
                    "fallback_ready": bool(rendered_pages),
                    "host_review_required": bool(unresolved_pages and rendered_pages),
                }
            else:
                stages["scanned_pdf_text"] = {"status": "not_applicable"}

            status = aggregate_status(stage["status"] for stage in stages.values())
            host_review_required = [
                name
                for name, stage in stages.items()
                if stage.get("host_review_required")
            ]
            content_path = item_dir / "content.txt"
            if normalized:
                atomic_write_text(content_path, normalized + "\n")

            chunks = natural_chunks(normalized, args.chunk_chars) if normalized else []
            chunk_index: list[dict[str, Any]] = []
            cursor = 0
            for index, chunk in enumerate(chunks, start=1):
                chunk_path = item_dir / "chunks" / f"chunk-{index:04d}.txt"
                atomic_write_text(chunk_path, chunk + "\n")
                start = normalized.find(chunk[: min(80, len(chunk))], cursor)
                start = cursor if start < 0 else start
                end = start + len(chunk)
                chunk_index.append(
                    {
                        "id": index,
                        "start": start,
                        "end": end,
                        "characters": len(chunk),
                        "preview": " ".join(chunk[:180].split()),
                        "path": relative_artifact(chunk_path, output),
                    }
                )
                cursor = end
            atomic_write_json(item_dir / "chunk-index.json", chunk_index)

            analysis = {
                "schema_version": 2,
                "id": item_id,
                "status": status,
                "status_scope": "text_extraction_only",
                "extraction_status": status,
                "distillation_status": "evidence_ready",
                "source": {
                    "name": label,
                    "type": path.suffix.casefold(),
                    "fingerprint": sampled_fingerprint(path),
                },
                "extraction": extraction,
                "stages": stages,
                "host_review_required": host_review_required,
                "outline": outline,
                "style_metrics": style_metrics(normalized) if normalized else None,
                "samples": sample_text(normalized) if normalized else {},
                "artifacts": {
                    "content": relative_artifact(content_path, output) if normalized else None,
                    "chunk_index": relative_artifact(item_dir / "chunk-index.json", output),
                    "chunk_count": len(chunks),
                    "rendered_pages": [
                        {
                            "page": page["page"],
                            "file": relative_artifact(item_dir / page["file"], output),
                        }
                        for page in rendered_pages
                    ],
                },
                "limitations": [
                    "Statistics describe visible form; the agent must infer argument, voice, and rhetorical purpose from the text.",
                    *([f"Scanned-PDF OCR unavailable: {ocr_error}"] if ocr_error else []),
                ],
                "semantic_review_required": [
                    "read the opening, middle, ending, and every block needed to support a judgment",
                    "record structure, argument or teaching moves, voice, rhythm, rhetorical actions, variables, and exceptions in a semantic card",
                    "attach artifact and locator evidence to every proposed transferable pattern",
                ],
            }
            atomic_write_json(analysis_path, analysis)
            items.append(
                {
                    "id": item_id,
                    "source": label,
                    "status": status,
                    "extraction_status": status,
                    "distillation_status": "evidence_ready",
                    "analysis": relative_artifact(analysis_path, output),
                    "host_review_required": host_review_required,
                }
            )
        except Exception as exc:
            items.append(
                {
                    "id": item_id,
                    "source": label,
                    "status": "failed",
                    "error": bounded_error(exc, path, output),
                }
            )

    counts = status_counts(items)
    corpus_analysis = build_corpus_analysis(output, items)
    atomic_write_json(output / "corpus-analysis.json", corpus_analysis)
    manifest = {
        "schema_version": 2,
        "kind": "document-analysis",
        "generated_at": utc_now(),
        "local_only": True,
        "status_scope": "text_extraction_only",
        "distillation_status": "evidence_ready",
        "source_count": len(sources),
        **counts,
        "scanned_pdf_ocr": {
            "mode": args.ocr_scanned,
            "available": ocr_engine is not None,
            "error": ocr_error,
        },
        "items": items,
        "corpus_analysis": "corpus-analysis.json",
    }
    atomic_write_json(output / "manifest.json", manifest)
    print(json.dumps(manifest, ensure_ascii=False, indent=2))
    return 1 if counts["failed_count"] else 0


if __name__ == "__main__":
    raise SystemExit(main())
